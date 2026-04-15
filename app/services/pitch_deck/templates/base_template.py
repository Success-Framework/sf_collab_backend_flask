"""
base_template.py
Shared python-pptx helpers used by all 4 pitch deck templates.
All measurement helpers, color utilities, and reusable shape builders live here.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy


# ---------------------------------------------------------------------------
# Slide dimensions — standard 16:9 widescreen
# ---------------------------------------------------------------------------
SLIDE_WIDTH  = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#RRGGBB' or 'RRGGBB' to RGBColor."""
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Shape / background builders
# ---------------------------------------------------------------------------

def fill_slide_background(slide, hex_color: str):
    """Fill the entire slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def add_rect(slide, left, top, width, height, fill_hex: str, transparency: float = 0.0):
    """
    Add a solid-filled rectangle shape.
    transparency: 0.0 = opaque, 1.0 = fully transparent
    """
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()  # no border
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(fill_hex)
    return shape


def add_circle(slide, left, top, diameter, fill_hex: str):
    """Add a solid-filled circle (oval with equal width/height)."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        9,  # MSO_AUTO_SHAPE_TYPE.OVAL
        left, top, diameter, diameter
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(fill_hex)
    return shape


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

def add_textbox(
    slide,
    text: str,
    left, top, width, height,
    font_name: str = "Calibri",
    font_size: int = 18,
    bold: bool = False,
    italic: bool = False,
    color_hex: str = "FFFFFF",
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
    word_wrap: bool = True
):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text

    font = run.font
    font.name = font_size and font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = hex_to_rgb(color_hex)

    return txBox


def add_bullet_textbox(
    slide,
    bullets: list,
    left, top, width, height,
    font_name: str = "Calibri",
    font_size: int = 14,
    color_hex: str = "FFFFFF",
    bullet_char: str = "•",
    line_spacing_pt: float = 6.0
):
    """
    Add a text box with bullet point lines.
    Each bullet is a separate paragraph for clean spacing.
    """
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if not bullet_text:
            continue

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT

        # Space before each bullet
        p.space_before = Pt(line_spacing_pt)

        run = p.add_run()
        run.text = f"{bullet_char}  {bullet_text}"

        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        font.color.rgb = hex_to_rgb(color_hex)

    return txBox


def add_label_value_pair(
    slide,
    label: str,
    value: str,
    left, top, width,
    label_color: str = "AAAAAA",
    value_color: str = "FFFFFF",
    font_name: str = "Calibri",
    label_size: int = 11,
    value_size: int = 22
):
    """
    Stacked label (small, muted) over value (large, bold).
    Used for stat callouts like TAM, SAM, SOM or metric boxes.
    """
    # Label
    add_textbox(
        slide, label.upper(),
        left, top, width, Inches(0.3),
        font_name=font_name, font_size=label_size,
        color_hex=label_color, bold=False
    )
    # Value
    add_textbox(
        slide, value,
        left, top + Inches(0.3), width, Inches(0.5),
        font_name=font_name, font_size=value_size,
        color_hex=value_color, bold=True
    )


# ---------------------------------------------------------------------------
# Slide factory
# ---------------------------------------------------------------------------

def new_blank_presentation() -> Presentation:
    """Create a new 16:9 Presentation with no default layouts."""
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_blank_slide(prs: Presentation):
    """Add a truly blank slide (blank layout = index 6)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


# ---------------------------------------------------------------------------
# Accent bar (left side vertical bar — the design motif we carry across slides)
# ---------------------------------------------------------------------------

def add_left_accent_bar(slide, color_hex: str, width: Inches = Inches(0.18)):
    """
    Add a thin vertical accent bar on the left edge of the slide.
    This is the visual motif carried across all slides in all templates.
    """
    add_rect(
        slide,
        left=Inches(0), top=Inches(0),
        width=width, height=SLIDE_HEIGHT,
        fill_hex=color_hex
    )


def add_bottom_strip(slide, color_hex: str, height: Inches = Inches(0.12)):
    """Thin accent strip along the bottom of the slide."""
    add_rect(
        slide,
        left=Inches(0), top=SLIDE_HEIGHT - height,
        width=SLIDE_WIDTH, height=height,
        fill_hex=color_hex
    )


# ---------------------------------------------------------------------------
# Slide number
# ---------------------------------------------------------------------------

def add_slide_number(slide, number: int, total: int, color_hex: str = "666666"):
    """Small slide number in bottom right corner."""
    add_textbox(
        slide,
        f"{number} / {total}",
        left=Inches(12.0), top=Inches(7.1),
        width=Inches(1.2), height=Inches(0.3),
        font_size=9, color_hex=color_hex,
        align=PP_ALIGN.RIGHT
    )