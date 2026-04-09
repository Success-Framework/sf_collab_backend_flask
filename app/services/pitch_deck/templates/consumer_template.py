"""
consumer_template.py
Theme: Coral Energy — Coral (#F96167) + Bold Black (#1A1A2E) + Warm White (#FFF8F0)
Visual motif: Bold top header bar on dark slides, warm cards on light slides
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from .base_template import (
    new_blank_presentation, add_blank_slide,
    fill_slide_background, add_rect, add_circle, add_textbox,
    add_bullet_textbox, add_left_accent_bar, add_bottom_strip,
    add_slide_number, SLIDE_WIDTH, SLIDE_HEIGHT
)

BG_DARK    = "1A1A2E"   # bold dark navy-black
BG_MID     = "16213E"   # slightly lighter dark
BG_LIGHT   = "FFF8F0"   # warm white
ACCENT     = "F96167"   # coral
ACCENT2    = "F9E795"   # gold yellow
TEXT_LITE  = "FFFFFF"
TEXT_DARK  = "1A1A2E"
TEXT_MID   = "444466"
CARD_BG    = "FFE8E8"   # pale coral for cards

TITLE_FONT = "Arial Black"
BODY_FONT  = "Calibri"
TOTAL_SLIDES = 8


def _dark_header(slide, title, num):
    # Bold coral top bar instead of left bar
    add_rect(slide, left=Inches(0), top=Inches(0),
             width=SLIDE_WIDTH, height=Inches(0.08), fill_hex=ACCENT)
    add_left_accent_bar(slide, ACCENT2, width=Inches(0.15))
    add_textbox(slide, title,
        left=Inches(0.5), top=Inches(0.22),
        width=Inches(11.0), height=Inches(0.7),
        font_name=TITLE_FONT, font_size=26, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)
    add_slide_number(slide, num, TOTAL_SLIDES, "667799")


def _light_header(slide, title, num):
    add_rect(slide, left=Inches(0), top=Inches(0),
             width=SLIDE_WIDTH, height=Inches(0.08), fill_hex=ACCENT)
    add_left_accent_bar(slide, ACCENT, width=Inches(0.15))
    add_bottom_strip(slide, ACCENT2)
    add_textbox(slide, title,
        left=Inches(0.5), top=Inches(0.22),
        width=Inches(11.0), height=Inches(0.7),
        font_name=TITLE_FONT, font_size=26, bold=True,
        color_hex=TEXT_DARK, align=PP_ALIGN.LEFT)
    add_slide_number(slide, num, TOTAL_SLIDES, "AAAAAA")


def _slide_cover(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)

    # Bold coral left panel
    add_rect(slide, left=Inches(0), top=Inches(0),
             width=Inches(4.5), height=SLIDE_HEIGHT, fill_hex=ACCENT)

    # Cross-hatching lines on coral panel (decorative)
    for i in range(6):
        add_rect(slide, left=Inches(0), top=Inches(0.7 + i * 1.1),
                 width=Inches(4.5), height=Inches(0.04), fill_hex="E04A50")

    add_textbox(slide, content.get("company_name",""),
        left=Inches(5.0), top=Inches(1.4),
        width=Inches(8.0), height=Inches(1.6),
        font_name=TITLE_FONT, font_size=46, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, content.get("tagline",""),
        left=Inches(5.0), top=Inches(3.2),
        width=Inches(8.0), height=Inches(0.65),
        font_name=BODY_FONT, font_size=18, italic=True,
        color_hex=ACCENT2, align=PP_ALIGN.LEFT)

    add_rect(slide, left=Inches(5.0), top=Inches(4.05),
             width=Inches(3.5), height=Inches(0.05), fill_hex=ACCENT)

    add_textbox(slide, content.get("founder_name",""),
        left=Inches(5.0), top=Inches(4.25),
        width=Inches(7.0), height=Inches(0.45),
        font_name=BODY_FONT, font_size=16, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, content.get("subtitle","Investor Presentation"),
        left=Inches(5.0), top=Inches(4.78),
        width=Inches(7.0), height=Inches(0.4),
        font_name=BODY_FONT, font_size=13,
        color_hex=ACCENT2, align=PP_ALIGN.LEFT)

    # Company initial on coral panel
    add_textbox(slide, content.get("company_name","?")[:1].upper(),
        left=Inches(0.5), top=Inches(2.5),
        width=Inches(3.5), height=Inches(2.5),
        font_name=TITLE_FONT, font_size=120, bold=True,
        color_hex="E04A50", align=PP_ALIGN.CENTER)


def _slide_problem(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_header(slide, content.get("slide_title","The Problem"), 2)

    add_rect(slide, left=Inches(0.5), top=Inches(1.2),
             width=Inches(12.5), height=Inches(0.72), fill_hex=BG_DARK)
    add_textbox(slide, content.get("headline",""),
        left=Inches(0.7), top=Inches(1.27),
        width=Inches(12.0), height=Inches(0.6),
        font_name=BODY_FONT, font_size=16, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)

    bullets = (content.get("bullets") or [])[:3]
    card_w = Inches(3.9); gap = Inches(0.22)

    for i, b in enumerate(bullets):
        cx = Inches(0.5) + i * (card_w + gap)
        add_rect(slide, left=cx, top=Inches(2.15), width=card_w, height=Inches(3.8), fill_hex=CARD_BG)
        add_rect(slide, left=cx, top=Inches(2.15), width=card_w, height=Inches(0.06), fill_hex=ACCENT)
        add_circle(slide, left=cx+Inches(0.2), top=Inches(2.3),
                   diameter=Inches(0.55), fill_hex=ACCENT)
        add_textbox(slide, str(i+1),
            left=cx+Inches(0.2), top=Inches(2.32),
            width=Inches(0.55), height=Inches(0.5),
            font_name=TITLE_FONT, font_size=14, bold=True,
            color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, b,
            left=cx+Inches(0.15), top=Inches(3.1),
            width=card_w-Inches(0.3), height=Inches(2.6),
            font_name=BODY_FONT, font_size=13,
            color_hex=TEXT_MID, align=PP_ALIGN.LEFT)


def _slide_solution(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)
    _dark_header(slide, content.get("slide_title","Our Solution"), 3)

    add_textbox(slide, content.get("headline",""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=ACCENT2, align=PP_ALIGN.LEFT)

    bullets = content.get("bullets") or []
    add_bullet_textbox(slide, bullets,
        left=Inches(0.5), top=Inches(2.0),
        width=Inches(6.3), height=Inches(4.8),
        font_name=BODY_FONT, font_size=15,
        color_hex=TEXT_LITE, bullet_char="★", line_spacing_pt=12)

    add_rect(slide, left=Inches(7.1), top=Inches(2.0),
             width=Inches(5.9), height=Inches(4.8), fill_hex=ACCENT)
    add_textbox(slide, "❤",
        left=Inches(9.0), top=Inches(2.8),
        width=Inches(2.5), height=Inches(2.0),
        font_name=BODY_FONT, font_size=72, bold=True,
        color_hex=BG_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Loved by Users",
        left=Inches(7.3), top=Inches(5.2),
        width=Inches(5.5), height=Inches(0.8),
        font_name=TITLE_FONT, font_size=16, bold=True,
        color_hex=BG_DARK, align=PP_ALIGN.CENTER)


def _slide_market(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_header(slide, content.get("slide_title","Market Opportunity"), 4)

    add_textbox(slide, content.get("headline",""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=TEXT_MID, align=PP_ALIGN.LEFT)

    keys   = ["tam","sam","som"]
    colors = [BG_DARK, ACCENT, ACCENT2]
    txt_c  = [TEXT_LITE, TEXT_LITE, TEXT_DARK]
    widths = [Inches(5.2), Inches(3.8), Inches(3.0)]
    heights= [Inches(3.5), Inches(2.9), Inches(2.3)]
    tops   = [Inches(2.0), Inches(2.3), Inches(2.6)]
    lefts  = [Inches(0.5), Inches(5.9), Inches(9.9)]
    abbrv  = ["TAM","SAM","SOM"]

    for i, (abbr, key) in enumerate(zip(abbrv, keys)):
        add_rect(slide, left=lefts[i], top=tops[i],
                 width=widths[i], height=heights[i], fill_hex=colors[i])
        add_textbox(slide, abbr,
            left=lefts[i]+Inches(0.2), top=tops[i]+Inches(0.15),
            width=widths[i]-Inches(0.3), height=Inches(0.45),
            font_name=TITLE_FONT, font_size=13, bold=True,
            color_hex=txt_c[i], align=PP_ALIGN.LEFT)
        add_textbox(slide, content.get(key,"—"),
            left=lefts[i]+Inches(0.2), top=tops[i]+Inches(0.8),
            width=widths[i]-Inches(0.3), height=heights[i]-Inches(1.0),
            font_name=BODY_FONT, font_size=15, bold=True,
            color_hex=txt_c[i], align=PP_ALIGN.LEFT)


def _slide_product(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_header(slide, content.get("slide_title","Our Product"), 5)

    add_textbox(slide, content.get("headline",""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=TEXT_MID, align=PP_ALIGN.LEFT)

    features = (content.get("features") or [])[:3]
    feat_w = Inches(3.9); gap = Inches(0.22)

    for i, feat in enumerate(features):
        cx = Inches(0.5) + i * (feat_w + gap)
        add_rect(slide, left=cx, top=Inches(2.0), width=feat_w, height=Inches(3.8), fill_hex=CARD_BG)
        add_rect(slide, left=cx, top=Inches(2.0), width=feat_w, height=Inches(0.08), fill_hex=ACCENT)
        add_circle(slide, left=cx+feat_w/2-Inches(0.4), top=Inches(2.2),
                   diameter=Inches(0.8), fill_hex=ACCENT)
        add_textbox(slide, feat.get("icon_label","")[:2].upper(),
            left=cx+feat_w/2-Inches(0.4), top=Inches(2.24),
            width=Inches(0.8), height=Inches(0.72),
            font_name=TITLE_FONT, font_size=11, bold=True,
            color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, feat.get("icon_label",""),
            left=cx+Inches(0.15), top=Inches(3.2),
            width=feat_w-Inches(0.3), height=Inches(0.5),
            font_name=TITLE_FONT, font_size=14, bold=True,
            color_hex=TEXT_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, feat.get("description",""),
            left=cx+Inches(0.15), top=Inches(3.8),
            width=feat_w-Inches(0.3), height=Inches(1.7),
            font_name=BODY_FONT, font_size=13,
            color_hex=TEXT_MID, align=PP_ALIGN.LEFT)


def _slide_traction(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)
    _dark_header(slide, content.get("slide_title","Traction & Milestones"), 6)

    add_textbox(slide, content.get("headline",""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=ACCENT2, align=PP_ALIGN.LEFT)

    metrics = (content.get("metrics") or [])[:3]
    m_w = Inches(3.9); gap = Inches(0.22)

    for i, m in enumerate(metrics):
        cx = Inches(0.5) + i * (m_w + gap)
        add_rect(slide, left=cx, top=Inches(1.9), width=m_w, height=Inches(1.8), fill_hex=ACCENT)
        add_textbox(slide, m.get("value","—"),
            left=cx+Inches(0.15), top=Inches(2.0),
            width=m_w-Inches(0.3), height=Inches(0.9),
            font_name=TITLE_FONT, font_size=28, bold=True,
            color_hex=BG_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, m.get("label",""),
            left=cx+Inches(0.15), top=Inches(2.92),
            width=m_w-Inches(0.3), height=Inches(0.5),
            font_name=BODY_FONT, font_size=11,
            color_hex=BG_DARK, align=PP_ALIGN.CENTER)

    milestones = (content.get("milestones") or [])[:4]
    add_textbox(slide, "KEY MILESTONES",
        left=Inches(0.5), top=Inches(4.0),
        width=Inches(5.0), height=Inches(0.35),
        font_name=BODY_FONT, font_size=11, bold=True,
        color_hex=ACCENT, align=PP_ALIGN.LEFT)
    add_rect(slide, left=Inches(0.5), top=Inches(4.52),
             width=Inches(12.5), height=Inches(0.04), fill_hex=ACCENT2)

    m_w2 = Inches(12.5) / max(len(milestones), 1)
    for i, ms in enumerate(milestones):
        cx = Inches(0.5) + i * m_w2
        add_circle(slide, left=cx+m_w2/2-Inches(0.1), top=Inches(4.42),
                   diameter=Inches(0.2), fill_hex=ACCENT2)
        add_textbox(slide, ms,
            left=cx, top=Inches(4.7),
            width=m_w2, height=Inches(1.5),
            font_name=BODY_FONT, font_size=12,
            color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)


def _slide_team(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_header(slide, content.get("slide_title","Our Team"), 7)

    add_textbox(slide, content.get("headline",""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=TEXT_MID, align=PP_ALIGN.LEFT)

    members = (content.get("members") or [])[:4]
    card_w = Inches(12.5) / max(len(members), 1) - Inches(0.2)

    for i, m in enumerate(members):
        cx = Inches(0.5) + i * (card_w + Inches(0.2))
        add_rect(slide, left=cx, top=Inches(2.0), width=card_w, height=Inches(4.0), fill_hex=CARD_BG)
        add_rect(slide, left=cx, top=Inches(2.0), width=card_w, height=Inches(0.08), fill_hex=ACCENT)
        add_circle(slide, left=cx+card_w/2-Inches(0.5), top=Inches(2.2),
                   diameter=Inches(1.0), fill_hex=ACCENT)
        add_textbox(slide, m.get("name","?")[:1].upper(),
            left=cx+card_w/2-Inches(0.5), top=Inches(2.25),
            width=Inches(1.0), height=Inches(0.9),
            font_name=TITLE_FONT, font_size=22, bold=True,
            color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, m.get("name",""),
            left=cx+Inches(0.1), top=Inches(3.35),
            width=card_w-Inches(0.2), height=Inches(0.45),
            font_name=TITLE_FONT, font_size=13, bold=True,
            color_hex=TEXT_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, m.get("role",""),
            left=cx+Inches(0.1), top=Inches(3.82),
            width=card_w-Inches(0.2), height=Inches(0.35),
            font_name=BODY_FONT, font_size=11,
            color_hex=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, m.get("highlight",""),
            left=cx+Inches(0.1), top=Inches(4.3),
            width=card_w-Inches(0.2), height=Inches(1.4),
            font_name=BODY_FONT, font_size=11,
            color_hex=TEXT_MID, align=PP_ALIGN.CENTER)


def _slide_ask(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)
    _dark_header(slide, content.get("slide_title","The Ask"), 8)

    add_textbox(slide, content.get("headline",""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=ACCENT2, align=PP_ALIGN.LEFT)

    add_rect(slide, left=Inches(0.5), top=Inches(1.9),
             width=Inches(5.5), height=Inches(2.0), fill_hex=ACCENT)
    add_textbox(slide, content.get("amount",""),
        left=Inches(0.7), top=Inches(2.0),
        width=Inches(5.1), height=Inches(1.0),
        font_name=TITLE_FONT, font_size=34, bold=True,
        color_hex=BG_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, "FUNDING ASK",
        left=Inches(0.7), top=Inches(3.1),
        width=Inches(5.1), height=Inches(0.5),
        font_name=BODY_FONT, font_size=13, bold=True,
        color_hex=BG_DARK, align=PP_ALIGN.CENTER)

    add_textbox(slide, "USE OF FUNDS",
        left=Inches(6.3), top=Inches(1.9),
        width=Inches(6.5), height=Inches(0.4),
        font_name=BODY_FONT, font_size=11, bold=True,
        color_hex=ACCENT2, align=PP_ALIGN.LEFT)

    funds = (content.get("use_of_funds") or [])[:4]
    for i, f in enumerate(funds):
        fy = Inches(2.45) + i * Inches(0.85)
        add_rect(slide, left=Inches(6.3), top=fy, width=Inches(6.5), height=Inches(0.7), fill_hex=BG_MID)
        add_textbox(slide, f.get("percentage",""),
            left=Inches(6.4), top=fy+Inches(0.12),
            width=Inches(1.0), height=Inches(0.5),
            font_name=TITLE_FONT, font_size=18, bold=True,
            color_hex=ACCENT, align=PP_ALIGN.LEFT)
        add_textbox(slide, f"{f.get('area','')} — {f.get('description','')}",
            left=Inches(7.5), top=fy+Inches(0.15),
            width=Inches(5.1), height=Inches(0.45),
            font_name=BODY_FONT, font_size=12,
            color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, content.get("closing_line",""),
        left=Inches(0.5), top=Inches(6.6),
        width=Inches(12.5), height=Inches(0.6),
        font_name=TITLE_FONT, font_size=13, italic=True,
        color_hex=ACCENT2, align=PP_ALIGN.CENTER)


def build_consumer_deck(content: dict):
    prs = new_blank_presentation()
    _slide_cover(prs, content.get("cover", {}))
    _slide_problem(prs, content.get("problem", {}))
    _slide_solution(prs, content.get("solution", {}))
    _slide_market(prs, content.get("market", {}))
    _slide_product(prs, content.get("product", {}))
    _slide_traction(prs, content.get("traction", {}))
    _slide_team(prs, content.get("team", {}))
    _slide_ask(prs, content.get("ask", {}))
    return prs