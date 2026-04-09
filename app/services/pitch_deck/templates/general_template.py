"""
general_template.py
Theme: Midnight Executive — Navy (#1E2761) + Ice Blue (#CADCFC) + Gold (#F5A623)
Visual motif: Left accent bar (gold) on all content slides, dark cover/ask slides
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from .base_template import (
    new_blank_presentation, add_blank_slide,
    fill_slide_background, add_rect, add_textbox,
    add_bullet_textbox, add_label_value_pair,
    add_left_accent_bar, add_bottom_strip,
    add_slide_number, hex_to_rgb, SLIDE_WIDTH, SLIDE_HEIGHT
)

# ── Palette ──────────────────────────────────────────────────────────────────
BG_DARK   = "1E2761"   # navy  — cover, ask slides
BG_LIGHT  = "F4F6FB"   # near-white — content slides
ACCENT    = "F5A623"   # gold  — accent bar, highlights
TEXT_DARK = "1E2761"   # navy  — on light slides
TEXT_MID  = "4A5580"   # slate — body on light slides
TEXT_LITE = "FFFFFF"   # white — on dark slides
ICE_BLUE  = "CADCFC"   # ice   — sub-headers on dark slides
CARD_BG   = "E8EDF8"   # pale blue card backgrounds

TITLE_FONT = "Georgia"
BODY_FONT  = "Calibri"

TOTAL_SLIDES = 8


# ── Helpers ──────────────────────────────────────────────────────────────────

def _dark_slide_header(slide, title: str, slide_num: int):
    """Shared dark-slide header: accent bar + title."""
    add_left_accent_bar(slide, ACCENT)
    add_textbox(
        slide, title,
        left=Inches(0.5), top=Inches(0.3),
        width=Inches(12.5), height=Inches(0.7),
        font_name=TITLE_FONT, font_size=32, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT
    )
    add_slide_number(slide, slide_num, TOTAL_SLIDES, color_hex="8899BB")


def _light_slide_header(slide, title: str, slide_num: int):
    """Shared light-slide header: accent bar + title + bottom strip."""
    add_left_accent_bar(slide, ACCENT)
    add_bottom_strip(slide, ACCENT)
    add_textbox(
        slide, title,
        left=Inches(0.5), top=Inches(0.25),
        width=Inches(12.5), height=Inches(0.7),
        font_name=TITLE_FONT, font_size=32, bold=True,
        color_hex=TEXT_DARK, align=PP_ALIGN.LEFT
    )
    add_slide_number(slide, slide_num, TOTAL_SLIDES, color_hex="AAAAAA")


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_cover(prs, content: dict):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)

    # Gold diagonal block (top-right design element)
    add_rect(slide,
        left=Inches(9.5), top=Inches(0),
        width=Inches(3.83), height=Inches(2.5),
        fill_hex=ACCENT
    )

    # Company name — large, centered-ish
    add_textbox(
        slide, content.get("company_name", ""),
        left=Inches(0.7), top=Inches(1.8),
        width=Inches(9.0), height=Inches(1.4),
        font_name=TITLE_FONT, font_size=52, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT
    )

    # Tagline
    add_textbox(
        slide, content.get("tagline", ""),
        left=Inches(0.7), top=Inches(3.3),
        width=Inches(9.0), height=Inches(0.6),
        font_name=BODY_FONT, font_size=20, bold=False, italic=True,
        color_hex=ICE_BLUE, align=PP_ALIGN.LEFT
    )

    # Divider line (thin rect)
    add_rect(slide,
        left=Inches(0.7), top=Inches(4.1),
        width=Inches(4.0), height=Inches(0.04),
        fill_hex=ACCENT
    )

    # Founder + subtitle
    add_textbox(
        slide, content.get("founder_name", ""),
        left=Inches(0.7), top=Inches(4.25),
        width=Inches(6.0), height=Inches(0.45),
        font_name=BODY_FONT, font_size=16, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT
    )
    add_textbox(
        slide, content.get("subtitle", "Investor Presentation"),
        left=Inches(0.7), top=Inches(4.75),
        width=Inches(6.0), height=Inches(0.4),
        font_name=BODY_FONT, font_size=13,
        color_hex=ICE_BLUE, align=PP_ALIGN.LEFT
    )


def _slide_problem(prs, content: dict):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_slide_header(slide, content.get("slide_title", "The Problem"), 2)

    # Headline box (dark card)
    add_rect(slide,
        left=Inches(0.5), top=Inches(1.2),
        width=Inches(12.5), height=Inches(0.75),
        fill_hex=BG_DARK
    )
    add_textbox(
        slide, content.get("headline", ""),
        left=Inches(0.7), top=Inches(1.28),
        width=Inches(12.0), height=Inches(0.6),
        font_name=BODY_FONT, font_size=17, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT
    )

    # 3 problem cards side by side
    bullets = (content.get("bullets") or [])[:3]
    card_w  = Inches(3.9)
    card_h  = Inches(3.5)
    gap     = Inches(0.22)
    start_x = Inches(0.5)
    card_y  = Inches(2.2)

    for i, bullet in enumerate(bullets):
        cx = start_x + i * (card_w + gap)
        add_rect(slide, left=cx, top=card_y, width=card_w, height=card_h, fill_hex=CARD_BG)
        # Number badge
        add_rect(slide, left=cx + Inches(0.15), top=card_y + Inches(0.18),
                 width=Inches(0.42), height=Inches(0.42), fill_hex=ACCENT)
        add_textbox(
            slide, str(i + 1),
            left=cx + Inches(0.15), top=card_y + Inches(0.16),
            width=Inches(0.42), height=Inches(0.42),
            font_name=BODY_FONT, font_size=14, bold=True,
            color_hex=TEXT_LITE, align=PP_ALIGN.CENTER
        )
        # Bullet text
        add_textbox(
            slide, bullet,
            left=cx + Inches(0.15), top=card_y + Inches(0.8),
            width=card_w - Inches(0.3), height=card_h - Inches(1.0),
            font_name=BODY_FONT, font_size=14,
            color_hex=TEXT_MID, align=PP_ALIGN.LEFT
        )


def _slide_solution(prs, content: dict):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_slide_header(slide, content.get("slide_title", "Our Solution"), 3)

    # Headline
    add_rect(slide,
        left=Inches(0.5), top=Inches(1.2),
        width=Inches(12.5), height=Inches(0.75),
        fill_hex=ACCENT
    )
    add_textbox(
        slide, content.get("headline", ""),
        left=Inches(0.7), top=Inches(1.28),
        width=Inches(12.0), height=Inches(0.6),
        font_name=BODY_FONT, font_size=17, bold=True,
        color_hex="1E2761", align=PP_ALIGN.LEFT
    )

    # Two-column: bullets left, visual right
    bullets = content.get("bullets") or []
    add_bullet_textbox(
        slide, bullets,
        left=Inches(0.5), top=Inches(2.2),
        width=Inches(6.5), height=Inches(4.5),
        font_name=BODY_FONT, font_size=15,
        color_hex=TEXT_DARK, bullet_char="→", line_spacing_pt=10
    )

    # Right: decorative card with icon area
    add_rect(slide,
        left=Inches(7.3), top=Inches(2.2),
        width=Inches(5.7), height=Inches(4.5),
        fill_hex=BG_DARK
    )
    add_textbox(
        slide, "✓",
        left=Inches(9.0), top=Inches(3.0),
        width=Inches(2.5), height=Inches(2.0),
        font_name=BODY_FONT, font_size=72, bold=True,
        color_hex=ACCENT, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, "Problem Solved",
        left=Inches(7.5), top=Inches(5.2),
        width=Inches(5.3), height=Inches(0.8),
        font_name=TITLE_FONT, font_size=18, bold=True,
        color_hex=ICE_BLUE, align=PP_ALIGN.CENTER
    )


def _slide_market(prs, content: dict):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)
    _dark_slide_header(slide, content.get("slide_title", "Market Opportunity"), 4)

    # Headline
    add_textbox(
        slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=ICE_BLUE, align=PP_ALIGN.LEFT
    )

    # TAM / SAM / SOM cards
    labels  = ["TAM", "SAM", "SOM"]
    keys    = ["tam", "sam", "som"]
    widths  = [Inches(5.5), Inches(4.0), Inches(3.0)]
    heights = [Inches(3.5), Inches(3.0), Inches(2.5)]
    tops    = [Inches(2.0), Inches(2.25), Inches(2.5)]
    lefts   = [Inches(0.5), Inches(6.2), Inches(10.4)]
    colors  = [ACCENT, "2D3A8C", "3D4EA0"]

    for i, (label, key) in enumerate(zip(labels, keys)):
        add_rect(slide,
            left=lefts[i], top=tops[i],
            width=widths[i], height=heights[i],
            fill_hex=colors[i]
        )
        add_textbox(
            slide, label,
            left=lefts[i] + Inches(0.2), top=tops[i] + Inches(0.2),
            width=widths[i] - Inches(0.3), height=Inches(0.5),
            font_name=TITLE_FONT, font_size=14, bold=True,
            color_hex=TEXT_LITE if i > 0 else BG_DARK,
            align=PP_ALIGN.LEFT
        )
        add_textbox(
            slide, content.get(key, "—"),
            left=lefts[i] + Inches(0.2), top=tops[i] + Inches(0.9),
            width=widths[i] - Inches(0.3), height=heights[i] - Inches(1.2),
            font_name=BODY_FONT, font_size=16, bold=True,
            color_hex=TEXT_LITE if i > 0 else BG_DARK,
            align=PP_ALIGN.LEFT
        )


def _slide_product(prs, content: dict):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_slide_header(slide, content.get("slide_title", "Our Product"), 5)

    add_textbox(
        slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=TEXT_MID, align=PP_ALIGN.LEFT
    )

    features = (content.get("features") or [])[:3]
    feat_w = Inches(3.9)
    feat_h = Inches(3.8)
    gap    = Inches(0.22)

    for i, feat in enumerate(features):
        cx = Inches(0.5) + i * (feat_w + gap)
        add_rect(slide, left=cx, top=Inches(2.0), width=feat_w, height=feat_h, fill_hex=BG_DARK)

        # Icon circle
        from .base_template import add_circle
        add_circle(slide, left=cx + Inches(1.5), top=Inches(2.2), diameter=Inches(0.7), fill_hex=ACCENT)
        add_textbox(
            slide, feat.get("icon_label", "")[:2].upper(),
            left=cx + Inches(1.5), top=Inches(2.22),
            width=Inches(0.7), height=Inches(0.66),
            font_name=BODY_FONT, font_size=11, bold=True,
            color_hex=BG_DARK, align=PP_ALIGN.CENTER
        )

        # Feature label
        add_textbox(
            slide, feat.get("icon_label", ""),
            left=cx + Inches(0.15), top=Inches(3.1),
            width=feat_w - Inches(0.3), height=Inches(0.5),
            font_name=TITLE_FONT, font_size=15, bold=True,
            color_hex=TEXT_LITE, align=PP_ALIGN.CENTER
        )

        # Feature description
        add_textbox(
            slide, feat.get("description", ""),
            left=cx + Inches(0.15), top=Inches(3.7),
            width=feat_w - Inches(0.3), height=Inches(1.8),
            font_name=BODY_FONT, font_size=13,
            color_hex=ICE_BLUE, align=PP_ALIGN.LEFT
        )


def _slide_traction(prs, content: dict):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_slide_header(slide, content.get("slide_title", "Traction & Milestones"), 6)

    add_textbox(
        slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=TEXT_MID, align=PP_ALIGN.LEFT
    )

    # Metrics row — large stat callouts
    metrics = (content.get("metrics") or [])[:3]
    m_w = Inches(3.9)
    m_h = Inches(1.8)
    gap  = Inches(0.22)

    for i, metric in enumerate(metrics):
        cx = Inches(0.5) + i * (m_w + gap)
        add_rect(slide, left=cx, top=Inches(1.9), width=m_w, height=m_h, fill_hex=BG_DARK)
        add_textbox(
            slide, metric.get("value", "—"),
            left=cx + Inches(0.15), top=Inches(2.0),
            width=m_w - Inches(0.3), height=Inches(0.9),
            font_name=TITLE_FONT, font_size=32, bold=True,
            color_hex=ACCENT, align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, metric.get("label", ""),
            left=cx + Inches(0.15), top=Inches(2.95),
            width=m_w - Inches(0.3), height=Inches(0.5),
            font_name=BODY_FONT, font_size=12,
            color_hex=ICE_BLUE, align=PP_ALIGN.CENTER
        )

    # Milestones timeline
    milestones = (content.get("milestones") or [])[:4]
    add_textbox(
        slide, "KEY MILESTONES",
        left=Inches(0.5), top=Inches(4.0),
        width=Inches(4.0), height=Inches(0.35),
        font_name=BODY_FONT, font_size=11, bold=True,
        color_hex=ACCENT, align=PP_ALIGN.LEFT
    )

    # Horizontal line
    add_rect(slide,
        left=Inches(0.5), top=Inches(4.55),
        width=Inches(12.5), height=Inches(0.04),
        fill_hex=ACCENT
    )

    m_w2 = Inches(12.5) / max(len(milestones), 1)
    for i, ms in enumerate(milestones):
        cx = Inches(0.5) + i * m_w2
        # Dot
        add_rect(slide, left=cx + m_w2/2 - Inches(0.1), top=Inches(4.44),
                 width=Inches(0.2), height=Inches(0.2), fill_hex=ACCENT)
        add_textbox(
            slide, ms,
            left=cx, top=Inches(4.7),
            width=m_w2, height=Inches(1.5),
            font_name=BODY_FONT, font_size=12,
            color_hex=TEXT_DARK, align=PP_ALIGN.CENTER
        )


def _slide_team(prs, content: dict):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_slide_header(slide, content.get("slide_title", "Our Team"), 7)

    add_textbox(
        slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=TEXT_MID, align=PP_ALIGN.LEFT
    )

    members = (content.get("members") or [])[:4]
    card_w = Inches(12.5) / max(len(members), 1) - Inches(0.2)
    card_h = Inches(4.0)

    for i, member in enumerate(members):
        cx = Inches(0.5) + i * (card_w + Inches(0.2))
        add_rect(slide, left=cx, top=Inches(2.0), width=card_w, height=card_h, fill_hex=CARD_BG)

        # Avatar circle
        from .base_template import add_circle
        add_circle(slide, left=cx + card_w/2 - Inches(0.5), top=Inches(2.2),
                   diameter=Inches(1.0), fill_hex=BG_DARK)
        add_textbox(
            slide, member.get("name", "?")[:1].upper(),
            left=cx + card_w/2 - Inches(0.5), top=Inches(2.25),
            width=Inches(1.0), height=Inches(0.9),
            font_name=TITLE_FONT, font_size=22, bold=True,
            color_hex=ACCENT, align=PP_ALIGN.CENTER
        )

        add_textbox(
            slide, member.get("name", ""),
            left=cx + Inches(0.1), top=Inches(3.35),
            width=card_w - Inches(0.2), height=Inches(0.45),
            font_name=TITLE_FONT, font_size=14, bold=True,
            color_hex=TEXT_DARK, align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, member.get("role", ""),
            left=cx + Inches(0.1), top=Inches(3.82),
            width=card_w - Inches(0.2), height=Inches(0.35),
            font_name=BODY_FONT, font_size=11,
            color_hex=ACCENT, align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, member.get("highlight", ""),
            left=cx + Inches(0.1), top=Inches(4.3),
            width=card_w - Inches(0.2), height=Inches(1.4),
            font_name=BODY_FONT, font_size=11,
            color_hex=TEXT_MID, align=PP_ALIGN.CENTER
        )


def _slide_ask(prs, content: dict):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)
    _dark_slide_header(slide, content.get("slide_title", "The Ask"), 8)

    add_textbox(
        slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=ICE_BLUE, align=PP_ALIGN.LEFT
    )

    # Funding amount — big callout
    add_rect(slide,
        left=Inches(0.5), top=Inches(1.9),
        width=Inches(5.5), height=Inches(2.0),
        fill_hex=ACCENT
    )
    add_textbox(
        slide, content.get("amount", ""),
        left=Inches(0.7), top=Inches(2.0),
        width=Inches(5.1), height=Inches(1.0),
        font_name=TITLE_FONT, font_size=36, bold=True,
        color_hex=BG_DARK, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, "FUNDING ASK",
        left=Inches(0.7), top=Inches(3.1),
        width=Inches(5.1), height=Inches(0.5),
        font_name=BODY_FONT, font_size=13, bold=True,
        color_hex=BG_DARK, align=PP_ALIGN.CENTER
    )

    # Use of funds breakdown
    add_textbox(
        slide, "USE OF FUNDS",
        left=Inches(6.3), top=Inches(1.9),
        width=Inches(6.5), height=Inches(0.4),
        font_name=BODY_FONT, font_size=11, bold=True,
        color_hex=ACCENT, align=PP_ALIGN.LEFT
    )

    funds = (content.get("use_of_funds") or [])[:4]
    for i, fund in enumerate(funds):
        fy = Inches(2.45) + i * Inches(0.85)
        add_rect(slide, left=Inches(6.3), top=fy, width=Inches(6.5), height=Inches(0.7), fill_hex="2D3A8C")
        add_textbox(
            slide, fund.get("percentage", ""),
            left=Inches(6.4), top=fy + Inches(0.12),
            width=Inches(1.0), height=Inches(0.5),
            font_name=TITLE_FONT, font_size=18, bold=True,
            color_hex=ACCENT, align=PP_ALIGN.LEFT
        )
        add_textbox(
            slide, f"{fund.get('area', '')} — {fund.get('description', '')}",
            left=Inches(7.5), top=fy + Inches(0.15),
            width=Inches(5.1), height=Inches(0.45),
            font_name=BODY_FONT, font_size=12,
            color_hex=TEXT_LITE, align=PP_ALIGN.LEFT
        )

    # Closing line
    add_textbox(
        slide, content.get("closing_line", ""),
        left=Inches(0.5), top=Inches(6.6),
        width=Inches(12.5), height=Inches(0.6),
        font_name=TITLE_FONT, font_size=14, italic=True,
        color_hex=ICE_BLUE, align=PP_ALIGN.CENTER
    )


# ── Public entry point ────────────────────────────────────────────────────────

def build_general_deck(content: dict):
    """
    Build the full General (Navy+Gold) pitch deck.
    Args:
        content: structured dict returned by content_generator.generate_slide_content()
    Returns:
        pptx.Presentation object
    """
    prs = new_blank_presentation()

    _slide_cover(prs, content.get("cover", {}))
    _slide_problem(prs, content.get("problem", {}))
    _slide_solution(prs, content.get("solution", {}))
    _slide_market(prs, content.get("market", {}))
    _slide_product(prs, content.get("product", {}))
    _slide_traction(prs, content.get("traction", {}))
    _slide_team(prs, content.get("team", {}))
    _slide_ask(prs, content.get("ask", {}))

    return prs