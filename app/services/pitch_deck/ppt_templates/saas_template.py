"""
saas_template.py
Theme: Deep Purple (#2D1B69) + Electric Blue (#4FC3F7) + White
Visual motif: Left accent bar (electric blue) + geometric top-right corner block
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from .base_template import (
    new_blank_presentation, add_blank_slide,
    fill_slide_background, add_rect, add_circle, add_textbox,
    add_bullet_textbox, add_left_accent_bar, add_bottom_strip,
    add_slide_number, SLIDE_WIDTH, SLIDE_HEIGHT
)

# ── Palette ───────────────────────────────────────────────────────────────────
BG_DARK    = "1A0A3C"   # deep dark purple — cover, dark slides
BG_MID     = "2D1B69"   # mid purple — cards, highlights
BG_LIGHT   = "F0F4FF"   # near-white lavender — content slides
ACCENT     = "4FC3F7"   # electric blue
ACCENT2    = "7C4DFF"   # violet — secondary accent
TEXT_LITE  = "FFFFFF"
TEXT_DARK  = "1A0A3C"
TEXT_MID   = "4A3880"
CARD_BG    = "E8EEFF"

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"

TOTAL_SLIDES = 8


def _dark_header(slide, title: str, num: int):
    add_left_accent_bar(slide, ACCENT)
    add_textbox(slide, title,
        left=Inches(0.5), top=Inches(0.28),
        width=Inches(11.0), height=Inches(0.7),
        font_name=TITLE_FONT, font_size=30, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)
    add_slide_number(slide, num, TOTAL_SLIDES, "5577AA")


def _light_header(slide, title: str, num: int):
    add_left_accent_bar(slide, ACCENT)
    add_bottom_strip(slide, ACCENT2)
    add_textbox(slide, title,
        left=Inches(0.5), top=Inches(0.28),
        width=Inches(11.0), height=Inches(0.7),
        font_name=TITLE_FONT, font_size=30, bold=True,
        color_hex=TEXT_DARK, align=PP_ALIGN.LEFT)
    add_slide_number(slide, num, TOTAL_SLIDES, "AAAAAA")


def _corner_geo(slide, color=ACCENT2, size=Inches(2.8)):
    """Geometric corner shape — top right design element."""
    add_rect(slide,
        left=SLIDE_WIDTH - size, top=Inches(0),
        width=size, height=size,
        fill_hex=color)


# ── Slides ────────────────────────────────────────────────────────────────────

def _slide_cover(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)

    _corner_geo(slide, ACCENT2, Inches(3.5))
    # Inner corner overlap
    _corner_geo(slide, ACCENT, Inches(2.2))

    add_textbox(slide, content.get("company_name", ""),
        left=Inches(0.7), top=Inches(1.6),
        width=Inches(9.5), height=Inches(1.4),
        font_name=TITLE_FONT, font_size=54, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, content.get("tagline", ""),
        left=Inches(0.7), top=Inches(3.15),
        width=Inches(9.0), height=Inches(0.6),
        font_name=BODY_FONT, font_size=19, italic=True,
        color_hex=ACCENT, align=PP_ALIGN.LEFT)

    add_rect(slide,
        left=Inches(0.7), top=Inches(3.95),
        width=Inches(3.5), height=Inches(0.04),
        fill_hex=ACCENT2)

    add_textbox(slide, content.get("founder_name", ""),
        left=Inches(0.7), top=Inches(4.1),
        width=Inches(6.0), height=Inches(0.45),
        font_name=BODY_FONT, font_size=16, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, content.get("subtitle", "Investor Presentation"),
        left=Inches(0.7), top=Inches(4.62),
        width=Inches(6.0), height=Inches(0.4),
        font_name=BODY_FONT, font_size=13,
        color_hex=ACCENT, align=PP_ALIGN.LEFT)


def _slide_problem(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_header(slide, content.get("slide_title", "The Problem"), 2)

    # Headline banner
    add_rect(slide, left=Inches(0.5), top=Inches(1.2),
             width=Inches(12.5), height=Inches(0.7), fill_hex=BG_MID)
    add_textbox(slide, content.get("headline", ""),
        left=Inches(0.7), top=Inches(1.27),
        width=Inches(12.0), height=Inches(0.6),
        font_name=BODY_FONT, font_size=16, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)

    bullets = (content.get("bullets") or [])[:3]
    card_w = Inches(3.9); gap = Inches(0.22)

    for i, b in enumerate(bullets):
        cx = Inches(0.5) + i * (card_w + gap)
        add_rect(slide, left=cx, top=Inches(2.15), width=card_w, height=Inches(3.8), fill_hex=CARD_BG)
        add_circle(slide, left=cx + Inches(0.15), top=Inches(2.3),
                   diameter=Inches(0.5), fill_hex=ACCENT)
        add_textbox(slide, str(i+1),
            left=cx + Inches(0.15), top=Inches(2.3),
            width=Inches(0.5), height=Inches(0.5),
            font_name=BODY_FONT, font_size=13, bold=True,
            color_hex=BG_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, b,
            left=cx + Inches(0.15), top=Inches(3.0),
            width=card_w - Inches(0.3), height=Inches(2.7),
            font_name=BODY_FONT, font_size=13,
            color_hex=TEXT_MID, align=PP_ALIGN.LEFT)


def _slide_solution(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)
    _dark_header(slide, content.get("slide_title", "Our Solution"), 3)

    add_textbox(slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=ACCENT, align=PP_ALIGN.LEFT)

    bullets = content.get("bullets") or []
    add_bullet_textbox(slide, bullets,
        left=Inches(0.5), top=Inches(2.0),
        width=Inches(6.3), height=Inches(4.8),
        font_name=BODY_FONT, font_size=15,
        color_hex=TEXT_LITE, bullet_char="▸", line_spacing_pt=12)

    # Right panel
    add_rect(slide, left=Inches(7.1), top=Inches(2.0),
             width=Inches(5.9), height=Inches(4.8), fill_hex=BG_MID)
    add_textbox(slide, "⚡",
        left=Inches(9.0), top=Inches(2.8),
        width=Inches(2.5), height=Inches(2.0),
        font_name=BODY_FONT, font_size=64, bold=True,
        color_hex=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Built for Scale",
        left=Inches(7.3), top=Inches(5.2),
        width=Inches(5.5), height=Inches(0.8),
        font_name=TITLE_FONT, font_size=17, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)


def _slide_market(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_header(slide, content.get("slide_title", "Market Opportunity"), 4)

    add_textbox(slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=TEXT_MID, align=PP_ALIGN.LEFT)

    labels = ["TAM", "SAM", "SOM"]
    keys   = ["tam", "sam", "som"]
    colors = [BG_MID, ACCENT2, ACCENT]
    txt_c  = [TEXT_LITE, TEXT_LITE, BG_DARK]
    widths = [Inches(5.2), Inches(3.8), Inches(3.0)]
    tops   = [Inches(2.0), Inches(2.3), Inches(2.6)]
    lefts  = [Inches(0.5), Inches(5.9), Inches(9.9)]
    heights = [Inches(3.5), Inches(2.9), Inches(2.3)]

    for i, (lbl, key) in enumerate(zip(labels, keys)):
        add_rect(slide, left=lefts[i], top=tops[i],
                 width=widths[i], height=heights[i], fill_hex=colors[i])
        add_textbox(slide, lbl,
            left=lefts[i]+Inches(0.2), top=tops[i]+Inches(0.15),
            width=widths[i]-Inches(0.3), height=Inches(0.45),
            font_name=TITLE_FONT, font_size=13, bold=True,
            color_hex=txt_c[i], align=PP_ALIGN.LEFT)
        add_textbox(slide, content.get(key, "—"),
            left=lefts[i]+Inches(0.2), top=tops[i]+Inches(0.8),
            width=widths[i]-Inches(0.3), height=heights[i]-Inches(1.0),
            font_name=BODY_FONT, font_size=15, bold=True,
            color_hex=txt_c[i], align=PP_ALIGN.LEFT)


def _slide_product(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)
    _dark_header(slide, content.get("slide_title", "Our Product"), 5)

    add_textbox(slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=ACCENT, align=PP_ALIGN.LEFT)

    features = (content.get("features") or [])[:3]
    feat_w = Inches(3.9); gap = Inches(0.22)

    for i, feat in enumerate(features):
        cx = Inches(0.5) + i * (feat_w + gap)
        add_rect(slide, left=cx, top=Inches(2.0), width=feat_w, height=Inches(3.8), fill_hex=BG_MID)
        add_circle(slide, left=cx+Inches(1.5), top=Inches(2.2),
                   diameter=Inches(0.7), fill_hex=ACCENT)
        add_textbox(slide, feat.get("icon_label", "")[:2].upper(),
            left=cx+Inches(1.5), top=Inches(2.22),
            width=Inches(0.7), height=Inches(0.66),
            font_name=BODY_FONT, font_size=11, bold=True,
            color_hex=BG_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, feat.get("icon_label", ""),
            left=cx+Inches(0.15), top=Inches(3.05),
            width=feat_w-Inches(0.3), height=Inches(0.5),
            font_name=TITLE_FONT, font_size=14, bold=True,
            color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, feat.get("description", ""),
            left=cx+Inches(0.15), top=Inches(3.65),
            width=feat_w-Inches(0.3), height=Inches(1.8),
            font_name=BODY_FONT, font_size=12,
            color_hex=ACCENT, align=PP_ALIGN.LEFT)


def _slide_traction(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_LIGHT)
    _light_header(slide, content.get("slide_title", "Traction & Milestones"), 6)

    add_textbox(slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=TEXT_MID, align=PP_ALIGN.LEFT)

    metrics = (content.get("metrics") or [])[:3]
    m_w = Inches(3.9); gap = Inches(0.22)

    for i, m in enumerate(metrics):
        cx = Inches(0.5) + i * (m_w + gap)
        add_rect(slide, left=cx, top=Inches(1.9), width=m_w, height=Inches(1.8), fill_hex=BG_MID)
        add_textbox(slide, m.get("value", "—"),
            left=cx+Inches(0.15), top=Inches(2.0),
            width=m_w-Inches(0.3), height=Inches(0.9),
            font_name=TITLE_FONT, font_size=30, bold=True,
            color_hex=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, m.get("label", ""),
            left=cx+Inches(0.15), top=Inches(2.92),
            width=m_w-Inches(0.3), height=Inches(0.5),
            font_name=BODY_FONT, font_size=11,
            color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)

    milestones = (content.get("milestones") or [])[:4]
    add_textbox(slide, "KEY MILESTONES",
        left=Inches(0.5), top=Inches(4.0),
        width=Inches(5.0), height=Inches(0.35),
        font_name=BODY_FONT, font_size=11, bold=True,
        color_hex=ACCENT2, align=PP_ALIGN.LEFT)

    add_rect(slide, left=Inches(0.5), top=Inches(4.52),
             width=Inches(12.5), height=Inches(0.04), fill_hex=ACCENT2)

    m_w2 = Inches(12.5) / max(len(milestones), 1)
    for i, ms in enumerate(milestones):
        cx = Inches(0.5) + i * m_w2
        add_circle(slide, left=cx+m_w2/2-Inches(0.1), top=Inches(4.42),
                   diameter=Inches(0.2), fill_hex=ACCENT2)
        add_textbox(slide, ms,
            left=cx, top=Inches(4.7),
            width=m_w2, height=Inches(1.5),
            font_name=BODY_FONT, font_size=12,
            color_hex=TEXT_DARK, align=PP_ALIGN.CENTER)


def _slide_team(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)
    _dark_header(slide, content.get("slide_title", "Our Team"), 7)

    add_textbox(slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=ACCENT, align=PP_ALIGN.LEFT)

    members = (content.get("members") or [])[:4]
    card_w = Inches(12.5) / max(len(members), 1) - Inches(0.2)

    for i, m in enumerate(members):
        cx = Inches(0.5) + i * (card_w + Inches(0.2))
        add_rect(slide, left=cx, top=Inches(2.0), width=card_w, height=Inches(4.0), fill_hex=BG_MID)
        add_circle(slide, left=cx+card_w/2-Inches(0.5), top=Inches(2.2),
                   diameter=Inches(1.0), fill_hex=ACCENT)
        add_textbox(slide, m.get("name","?")[:1].upper(),
            left=cx+card_w/2-Inches(0.5), top=Inches(2.25),
            width=Inches(1.0), height=Inches(0.9),
            font_name=TITLE_FONT, font_size=22, bold=True,
            color_hex=BG_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, m.get("name",""),
            left=cx+Inches(0.1), top=Inches(3.35),
            width=card_w-Inches(0.2), height=Inches(0.45),
            font_name=TITLE_FONT, font_size=13, bold=True,
            color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, m.get("role",""),
            left=cx+Inches(0.1), top=Inches(3.82),
            width=card_w-Inches(0.2), height=Inches(0.35),
            font_name=BODY_FONT, font_size=11,
            color_hex=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, m.get("highlight",""),
            left=cx+Inches(0.1), top=Inches(4.3),
            width=card_w-Inches(0.2), height=Inches(1.4),
            font_name=BODY_FONT, font_size=11,
            color_hex="BBBBDD", align=PP_ALIGN.CENTER)


def _slide_ask(prs, content):
    slide = add_blank_slide(prs)
    fill_slide_background(slide, BG_DARK)
    _dark_header(slide, content.get("slide_title", "The Ask"), 8)

    add_textbox(slide, content.get("headline", ""),
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(12.5), height=Inches(0.5),
        font_name=BODY_FONT, font_size=15, italic=True,
        color_hex=ACCENT, align=PP_ALIGN.LEFT)

    add_rect(slide, left=Inches(0.5), top=Inches(1.9),
             width=Inches(5.5), height=Inches(2.0), fill_hex=ACCENT2)
    add_textbox(slide, content.get("amount",""),
        left=Inches(0.7), top=Inches(2.0),
        width=Inches(5.1), height=Inches(1.0),
        font_name=TITLE_FONT, font_size=36, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "FUNDING ASK",
        left=Inches(0.7), top=Inches(3.1),
        width=Inches(5.1), height=Inches(0.5),
        font_name=BODY_FONT, font_size=13, bold=True,
        color_hex=TEXT_LITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "USE OF FUNDS",
        left=Inches(6.3), top=Inches(1.9),
        width=Inches(6.5), height=Inches(0.4),
        font_name=BODY_FONT, font_size=11, bold=True,
        color_hex=ACCENT, align=PP_ALIGN.LEFT)

    funds = (content.get("use_of_funds") or [])[:4]
    for i, f in enumerate(funds):
        fy = Inches(2.45) + i * Inches(0.85)
        add_rect(slide, left=Inches(6.3), top=fy, width=Inches(6.5), height=Inches(0.7), fill_hex=BG_MID)
        add_textbox(slide, f.get("percentage",""),
            left=Inches(6.4), top=fy+Inches(0.12),
            width=Inches(1.0), height=Inches(0.5),
            font_name=TITLE_FONT, font_size=18, bold=True,
            color_hex=ACCENT, align=PP_ALIGN.LEFT)
        add_textbox(slide, f"{f.get('area','')} — {f.get('description','')}",
            left=Inches(7.5), top=fy+Inches(0.15),
            width=Inches(5.1), height=Inches(0.45),
            font_name=BODY_FONT, font_size=12,
            color_hex=TEXT_LITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, content.get("closing_line",""),
        left=Inches(0.5), top=Inches(6.6),
        width=Inches(12.5), height=Inches(0.6),
        font_name=TITLE_FONT, font_size=14, italic=True,
        color_hex=ACCENT, align=PP_ALIGN.CENTER)


def build_saas_deck(content: dict):
    prs = new_blank_presentation()
    _slide_cover(prs, content.get("cover", {}))
    _slide_problem(prs, content.get("problem", {}))
    _slide_solution(prs, content.get("solution", {}))
    _slide_market(prs, content.get("market", {}))
    _slide_product(prs, content.get("product", {}))
    _slide_traction(prs, content.get("traction", {}))
    _slide_team(prs, content.get("team", {}))
    _slide_ask(prs, content.get("ask", {}))
    return prs