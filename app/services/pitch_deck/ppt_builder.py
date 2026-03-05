# app/services/pitch_deck/ppt_builder.py

import os
import re
from flask import current_app
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
# from app.services.pitch_deck.themes import get_theme

TEMPLATE_MAP = {
    "creative": "creative.pptx",
    "startup": "startup.pptx",
    "corporate": "corporate.pptx"
}



class PitchDeckPPTBuilder:

    def __init__(self, deck_json, theme_type, deck_id):
        self.deck_json = deck_json
        self.deck_id = deck_id

        base_upload_folder = current_app.config.get("UPLOAD_FOLDER")

        TEMPLATE_MAP = {
            "creative": "creative.pptx",
            "startup": "startup.pptx",
            "corporate": "corporate.pptx"
        }

# Use default if invalid
        template_filename = TEMPLATE_MAP.get(theme_type, "startup.pptx")

        template_path = os.path.join(
            base_upload_folder,
            "pitch_deck_templates",
            template_filename
        )

        if not os.path.exists(template_path):
            raise Exception(f"Template file not found at: {template_path}")

        self.prs = Presentation(template_path)


        self.output_folder = os.path.join(base_upload_folder, "pitch_decks")
        os.makedirs(self.output_folder, exist_ok=True)

    # --------------------------------------------------

    def build(self):

        for slide_data in self.deck_json["slides"]:
            self._add_slide(slide_data)

        filename = f"{self.deck_id}.pptx"
        filepath = os.path.join(self.output_folder, filename)
        self.prs.save(filepath)

        return filepath

    # --------------------------------------------------

    def _add_slide(self, slide_data):

        slide_type = slide_data.get("type")

        # Layout assumptions:
        # 0 = Cover
        # 1 = Title + Content

        if slide_type == "cover":
            slide_layout = self.prs.slide_layouts[0]
        else:
            slide_layout = self.prs.slide_layouts[1]

        slide = self.prs.slides.add_slide(slide_layout)

        # Cover Slide
        if slide_type == "cover":
            slide.placeholders[0].text = slide_data.get("title", "")

            bullets = slide_data.get("bullets", [])
            if len(bullets) > 0 and len(slide.placeholders) > 1:
                slide.placeholders[1].text = bullets[0]

            return

        # Title
        slide.placeholders[0].text = slide_data.get("title", "")

        # Financial Slide
        if slide_type == "financials":
            self._add_financial_chart(slide, slide_data.get("bullets", []))
        else:
            bullets = slide_data.get("bullets", [])[:3]
            slide.placeholders[1].text = "\n".join(bullets)

    # --------------------------------------------------

    def _add_financial_chart(self, slide, bullets):

        years = []
        values = []

        for bullet in bullets:
            match = re.search(r"(Year\s*\d+|\d{4}).*?\$([\d\.]+)", bullet, re.IGNORECASE)
            if match:
                years.append(match.group(1))
                values.append(float(match.group(2)))

        if len(values) < 3:
            years = ["Year 1", "Year 2", "Year 3"]
            values = [100, 300, 800]

        chart_data = ChartData()
        chart_data.categories = years[:3]
        chart_data.add_series("Revenue", values[:3])

        # Replace content placeholder with chart
        content = slide.placeholders[1]
        left = content.left
        top = content.top
        width = content.width
        height = content.height

        slide.shapes._spTree.remove(content._element)

        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            left,
            top,
            width,
            height,
            chart_data
        )
